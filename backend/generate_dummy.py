import os
from fpdf import FPDF
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt

templates_dir = "storage/app/public/templates"
os.makedirs(templates_dir, exist_ok=True)

# 1. Generate PDF
pdf = FPDF()
pdf.add_page()
pdf.set_font("Arial", size=15)
pdf.cell(200, 10, txt="Hasil Konversi Dokumen - SmartDoc AI", ln=1, align='C')
pdf.set_font("Arial", size=12)
pdf.cell(200, 10, txt="Ini adalah simulasi hasil konversi. Dokumen ini 100% valid dan bisa dibuka.", ln=1, align='C')
pdf.output(os.path.join(templates_dir, "converted.pdf"))

# 2. Generate DOCX
doc = Document()
doc.add_heading('Hasil Konversi Dokumen - SmartDoc AI', 0)
doc.add_paragraph('Ini adalah simulasi hasil konversi. Dokumen ini 100% valid dan bisa dibuka.')
doc.save(os.path.join(templates_dir, "converted.docx"))

# 3. Generate XLSX
wb = Workbook()
ws = wb.active
ws.title = "Hasil Konversi"
ws['A1'] = "Hasil Konversi Dokumen - SmartDoc AI"
ws['A2'] = "Ini adalah simulasi hasil konversi. Dokumen ini 100% valid dan bisa dibuka."
wb.save(os.path.join(templates_dir, "converted.xlsx"))

# 4. Generate PPTX
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hasil Konversi - SmartDoc AI"
subtitle.text = "Ini adalah simulasi hasil konversi.\nDokumen ini 100% valid dan bisa dibuka."
prs.save(os.path.join(templates_dir, "converted.pptx"))

print("Semua file dummy valid berhasil dibuat di", templates_dir)
